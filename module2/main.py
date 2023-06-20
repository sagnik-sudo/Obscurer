from fastapi import FastAPI, UploadFile, File
from pdf2image import convert_from_path
from PIL import Image
import textract
import docx2txt
import os
import magic
import re
import pandas as pd
from pptx import Presentation

app = FastAPI()


def convert_pdf_to_text(file_path):
    images = convert_from_path(file_path)
    text = ""
    for img in images:
        img_path = f"{file_path}.jpg"
        img.save(img_path)
        text += textract.process(img_path).decode("utf-8")
        os.remove(img_path)
    return text


def convert_audio_to_text(file_path):
    # Implement audio-to-text conversion logic using a library like SpeechRecognition
    # Return the extracted text
    pass


def parse_hl7(file_path):
    # Implement HL7 parsing logic to extract text
    # Return the extracted text
    pass


def process_text(text):
    # Replace "\n" and "\t" with actual line breaks and tabs
    text = re.sub(r"\\n", "\n", text)
    text = re.sub(r"\\t", "\t", text)
    return text


def read_txt(file_path):
    with open(file_path, "r") as f:
        text = f.read()
    return text


def read_csv(file_path):
    df = pd.read_csv(file_path)
    text = df.to_string(index=False)
    return text


def read_xlsx(file_path):
    df = pd.read_excel(file_path, engine="openpyxl")
    text = df.to_string(index=False)
    return text


# def read_pptx(file_path):
#     prs = python_pptx.Presentation(file_path)
#     text = ""
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text + "\n"
#     return text


@app.post("/convert")
async def convert_to_text(file: UploadFile = File(...)):
    file_extension = file.filename.split(".")[-1]
    file_path = f"uploads/{file.filename}"

    with open(file_path, "wb") as f:
        f.write(file.file.read())

    if file_extension.lower() == "pdf":
        # Check if the file is actually a PDF
        file_mime = magic.from_file(file_path, mime=True)
        if file_mime != "application/pdf":
            return {"error": "Invalid PDF file."}

        text = convert_pdf_to_text(file_path)
    elif file_extension.lower() in ["jpg", "jpeg", "png", "bmp", "gif"]:
        text = textract.process(file_path).decode("utf-8")
    elif file_extension.lower() == "docx":
        text = docx2txt.process(file_path)
    # elif file_extension.lower() == "hl7":
    #     text = parse_hl7(file_path)
    # elif file_extension.lower() in ["wav", "mp3"]:
    #     text = convert_audio_to_text(file_path)
    elif file_extension.lower() == "txt":
        text = read_txt(file_path)
    elif file_extension.lower() == "csv":
        text = read_csv(file_path)
    elif file_extension.lower() == "xlsx":
        text = read_xlsx(file_path)
    # elif file_extension.lower() == "pptx":
    #     text = read_pptx(file_path)
    else:
        return {"error": "Unsupported file format."}

    processed_text = process_text(text)
    os.remove(file_path)
    return {"text": processed_text}


# if __name__ == "__main__":
#     uvicorn.run(app, host="0.0.0.0", port=8000)
